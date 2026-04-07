"""
PPT Renderer
=============
PresentationDesign(JSON) + raw data → python-pptx PPTX 파일 생성.

LLM이 결정한 디자인(색상, 폰트, 배치)을 적용하고,
data_ref placeholder를 실제 테이블/차트로 치환합니다.
"""
from __future__ import annotations

import io
import logging
from datetime import date, datetime
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from common import PARA_COLUMNS, PT1C_COLUMNS, GMS_COLUMNS, HIGHER_IS_BETTER
from ppt_llm_designer import PresentationDesign, SlideDesign, ElementDesign

logger = logging.getLogger("yield_agent.ppt_renderer")

# matplotlib 한글 폰트
_KR_FONTS = ["Malgun Gothic", "NanumGothic", "AppleGothic", "DejaVu Sans"]
for _f in _KR_FONTS:
    try:
        matplotlib.font_manager.findfont(_f, fallback_to_default=False)
        plt.rcParams["font.family"] = _f
        break
    except Exception:
        continue
plt.rcParams["axes.unicode_minus"] = False

KEY_PT1H_PARAMS = ["VTH", "IDSAT", "IDLIN", "IOFF", "ION", "IGATE", "IDDQ",
                   "VMIN", "FMAX", "GM_MAX", "SS", "DIBL"]


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """'#RRGGBB' → RGBColor."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "1E293B"
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def render_presentation(prs: Presentation, design: PresentationDesign,
                        state: dict[str, Any]) -> None:
    """PresentationDesign에 따라 prs에 슬라이드를 추가한다."""
    cs = design.color_scheme
    font = design.font_family or "맑은 고딕"

    from ppt_builder import _get_layout, _clear_placeholders

    for slide_design in design.slides:
        purpose = "title" if slide_design.layout in ("title", "ending") else "blank"
        layout = _get_layout(prs, purpose)
        slide = prs.slides.add_slide(layout)
        _clear_placeholders(slide)

        # 배경색
        if slide_design.background_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(slide_design.background_color)

        # 슬라이드 헤더 (data_overview만)
        if slide_design.layout == "data_overview":
            _add_slide_header(slide, slide_design.title, cs, font)

        # 각 요소 렌더링
        for elem in slide_design.elements:
            _render_element(slide, elem, design, state, font)


def _add_slide_header(slide, text: str, cs, font: str):
    """슬라이드 상단 헤더 (제목 + 언더라인)."""
    _add_textbox(slide, text,
                 left=Inches(0.8), top=Inches(0.3),
                 width=Inches(11), height=Inches(0.7),
                 font_size=Pt(24), font_color=_hex_to_rgb(cs.text),
                 bold=True, font_name=font)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.05),
        width=Inches(2), height=Pt(3),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(cs.accent)
    bar.line.fill.background()


def _render_element(slide, elem: ElementDesign, design: PresentationDesign,
                    state: dict[str, Any], font: str):
    """data_ref에 따라 적절한 요소를 렌더링."""
    cs = design.color_scheme
    ref = elem.data_ref

    if ref == "title_text":
        lotcd = state.get("lotcd", "Unknown")
        _add_textbox(slide, f"{lotcd} 수율 분석 리포트",
                     left=Inches(elem.left), top=Inches(elem.top),
                     width=Inches(elem.width), height=Inches(elem.height),
                     font_size=Pt(elem.resolved_style.font_size or 40),
                     font_color=_hex_to_rgb(elem.resolved_style.font_color or cs.text),
                     bold=elem.resolved_style.bold, font_name=font,
                     alignment=PP_ALIGN.CENTER)

    elif ref == "subtitle_text":
        ref_date = state.get("ref_date", "")
        unit = state.get("unit", "weekly")
        unit_label = {"weekly": "주간", "monthly": "월간", "daily": "일별"}.get(unit, unit)
        try:
            dt = datetime.strptime(ref_date, "%Y%m%d")
            date_str = dt.strftime("%Y년 %m월 %d일")
        except ValueError:
            date_str = ref_date
        subtitle = f"기준일: {date_str}  |  단위: {unit_label}  |  생성: {date.today().strftime('%Y-%m-%d')}"
        _add_textbox(slide, subtitle,
                     left=Inches(elem.left), top=Inches(elem.top),
                     width=Inches(elem.width), height=Inches(elem.height),
                     font_size=Pt(elem.resolved_style.font_size or 16),
                     font_color=_hex_to_rgb(elem.resolved_style.font_color or "#CBD5E1"),
                     font_name=font, alignment=PP_ALIGN.CENTER)

    elif ref == "ending_text":
        _add_textbox(slide, "Thank You",
                     left=Inches(elem.left), top=Inches(elem.top),
                     width=Inches(elem.width), height=Inches(elem.height),
                     font_size=Pt(elem.resolved_style.font_size or 36),
                     font_color=_hex_to_rgb(elem.resolved_style.font_color or "#FFFFFF"),
                     bold=True, font_name=font, alignment=PP_ALIGN.CENTER)

    elif ref == "ending_footer":
        _add_textbox(slide, f"Generated by Yield Agent  |  {date.today().strftime('%Y-%m-%d')}",
                     left=Inches(elem.left), top=Inches(elem.top),
                     width=Inches(elem.width), height=Inches(elem.height),
                     font_size=Pt(elem.resolved_style.font_size or 14),
                     font_color=_hex_to_rgb(elem.resolved_style.font_color or "#CBD5E1"),
                     font_name=font, alignment=PP_ALIGN.CENTER)

    elif ref == "unified_table":
        _render_unified_table(slide, elem, design, state, font)

    elif ref.startswith("scatter_"):
        _render_scatter_chart(slide, elem, ref, design, state)

    elif ref == "anomaly_summary":
        anomaly_params = state.get("anomaly_params", [])
        if anomaly_params:
            lines = [f"{a['param']}: {a.get('direction','?')} ({a.get('change_pct',0):+.1f}%)"
                     for a in anomaly_params]
            text = "이상 파라미터: " + ", ".join(lines)
        else:
            text = "이상 파라미터 없음"
        _add_textbox(slide, text,
                     left=Inches(elem.left), top=Inches(elem.top),
                     width=Inches(elem.width), height=Inches(elem.height),
                     font_size=Pt(elem.resolved_style.font_size or 10),
                     font_color=_hex_to_rgb(elem.resolved_style.font_color or cs.text),
                     font_name=font)

    elif ref == "extra_content":
        # SlideDesign.extra_content에 저장된 LLM 요약 텍스트를 렌더링
        # slide_design은 외부에서 전달되지 않으므로, 렌더링 시점에
        # slide_design.extra_content를 state에 임시 주입하여 전달
        text = state.get("_current_extra_content", "")
        if text:
            _render_multiline_text(slide, text, elem, cs, font)


# ================================================================
# 셀 텍스트 90도 회전 (OOXML vert270 — Windows PowerPoint 지원)
# ================================================================
def _set_cell_vertical(cell):
    """셀 텍스트를 90도 회전 (bottom-to-top)."""
    _ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    txBody = cell._tc.find(f"{{{_ns}}}txBody")
    if txBody is not None:
        bodyPr = txBody.find(f"{{{_ns}}}bodyPr")
        if bodyPr is not None:
            bodyPr.set("vert", "vert270")


# ================================================================
# 통합 멀티헤더 테이블 (전체 파라미터, vert270 회전 헤더)
# ================================================================
def _render_unified_table(slide, elem: ElementDesign, design: PresentationDesign,
                          state: dict[str, Any], font: str):
    """PT1H/PT1C/GMS 통합 멀티헤더 테이블 — 전체 32 파라미터, vert270 회전 헤더."""
    cs = design.color_scheme
    weeks_data = state.get("weeks_data", [])
    anomaly_params = state.get("anomaly_params", [])
    unit = state.get("unit", "weekly")

    period_label = {"weekly": "WEEK", "monthly": "MONTH", "daily": "DATE"}.get(unit, "WEEK")
    anomaly_map = {a["param"]: a for a in anomaly_params} if anomaly_params else {}

    # 전체 파라미터 (25 + 2 + 5 = 32)
    groups = [
        ("PT1H", cs.table_header, [(c, c) for c in PARA_COLUMNS]),
        ("PT1C", "#6D28D9", [(c, f"pt1c_{c}") for c in PT1C_COLUMNS]),
        ("GMS", "#059669", [(c.upper(), f"gms_{c}") for c in GMS_COLUMNS]),
    ]

    all_params = []
    group_spans = []
    col_offset = 1
    for gname, gcolor, params in groups:
        group_spans.append((gname, gcolor, col_offset, len(params)))
        all_params.extend(params)
        col_offset += len(params)

    n_cols = 1 + len(all_params)
    n_rows = 2 + len(weeks_data)  # 그룹헤더 + 파라미터헤더(회전) + 데이터
    header_row_h = Inches(0.55)  # 회전 헤더 높이
    group_row_h = Inches(0.22)
    data_row_h = Inches(0.22)
    table_height = group_row_h + header_row_h + data_row_h * len(weeks_data)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(elem.left), Inches(elem.top),
        Inches(elem.width), table_height,
    )
    table = table_shape.table

    # 행 높이 설정
    table.rows[0].height = group_row_h
    table.rows[1].height = header_row_h
    for i in range(2, n_rows):
        table.rows[i].height = data_row_h

    # 컬럼 폭: period 넓게, 나머지 균등
    table.columns[0].width = Inches(0.7)
    param_col_w = int((Inches(elem.width) - Inches(0.7)) / len(all_params))
    for j in range(1, n_cols):
        table.columns[j].width = param_col_w

    # row 0: 그룹 헤더 (merge) — period 셀은 세로 merge
    table.cell(0, 0).merge(table.cell(1, 0))
    table.cell(0, 0).text = period_label
    _style_cell(table.cell(0, 0), Pt(6), True,
                RGBColor(0xFF, 0xFF, 0xFF), _hex_to_rgb(cs.table_header), font)

    for gname, gcolor, start, span in group_spans:
        if span > 1:
            table.cell(0, start).merge(table.cell(0, start + span - 1))
        table.cell(0, start).text = gname
        _style_cell(table.cell(0, start), Pt(6), True,
                    RGBColor(0xFF, 0xFF, 0xFF), _hex_to_rgb(gcolor), font)

    # row 1: 파라미터 이름 (vert270 회전 — Windows PowerPoint에서 작동)
    for j, (display, data_key) in enumerate(all_params):
        cell = table.cell(1, 1 + j)
        cell.text = display
        base_param = data_key.replace("pt1c_", "").replace("gms_", "")
        if base_param in anomaly_map:
            bg = _hex_to_rgb(cs.danger) if anomaly_map[base_param].get("direction") == "열화" else _hex_to_rgb(cs.success)
        else:
            bg = _hex_to_rgb("#64748B")
        _style_cell(cell, Pt(5), True, RGBColor(0xFF, 0xFF, 0xFF), bg, font)
        _set_cell_vertical(cell)

    # row 2+: 데이터
    last_idx = len(weeks_data) - 1
    alt_bg = _hex_to_rgb(cs.table_alt_row)
    white_bg = RGBColor(0xFF, 0xFF, 0xFF)
    text_color = _hex_to_rgb(cs.text)
    gray_color = _hex_to_rgb("#475569")

    for i, wd in enumerate(weeks_data):
        row_idx = i + 2
        bg = white_bg if i % 2 == 0 else alt_bg

        table.cell(row_idx, 0).text = str(wd.get("week", "?"))
        _style_cell(table.cell(row_idx, 0), Pt(5), False, text_color, bg, font)

        for j, (display, data_key) in enumerate(all_params):
            cell = table.cell(row_idx, 1 + j)
            cell.text = _fmt_val(wd.get(data_key, "-"))
            base_param = data_key.replace("pt1c_", "").replace("gms_", "")
            if i == last_idx and base_param in anomaly_map:
                a = anomaly_map[base_param]
                fc = _hex_to_rgb(cs.danger) if a.get("direction") == "열화" else _hex_to_rgb(cs.success)
                _style_cell(cell, Pt(5), True, fc, bg, font)
            else:
                _style_cell(cell, Pt(5), False, gray_color, bg, font)


# ================================================================
# Scatter Chart 렌더링
# ================================================================
def _render_scatter_chart(slide, elem: ElementDesign, ref: str,
                          design: PresentationDesign, state: dict[str, Any]):
    """scatter plot을 matplotlib로 생성하여 슬라이드에 삽입."""
    cs = design.color_scheme
    weeks_data = state.get("weeks_data", [])
    anomaly_params = state.get("anomaly_params", [])
    periods = [wd.get("week", f"P{i}") for i, wd in enumerate(weeks_data)]

    mpl_colors = [cs.accent, cs.danger, cs.success, "#F59E0B", "#8B5CF6", "#06B6D4"]

    fig, ax = plt.subplots(figsize=(elem.width * 0.95, elem.height * 0.85))
    fig.patch.set_facecolor("white")

    if ref == "scatter_pt1h_pt1c":
        params = []
        for a in anomaly_params:
            p = a["param"]
            if p in PARA_COLUMNS:
                params.append((p, p))
            elif p in PT1C_COLUMNS:
                params.append((p, f"pt1c_{p}"))
        _plot_multi_trend(ax, params[:6], weeks_data, periods, mpl_colors, "PT1H / PT1C")

    elif ref == "scatter_degradation":
        if anomaly_params:
            for idx, a in enumerate(anomaly_params[:6]):
                param = a["param"]
                data_key = f"pt1c_{param}" if param in PT1C_COLUMNS else param
                values = _extract_values(weeks_data, data_key)
                x_valid = [i for i, v in enumerate(values) if v is not None]
                y_valid = [values[i] for i in x_valid]
                color = cs.danger if a.get("direction") == "열화" else cs.success
                marker = "v" if a.get("direction") == "열화" else "^"
                direction = a.get("direction", "?")
                pct = a.get("change_pct", 0.0)
                label = f"{param} ({direction} {pct:+.1f}%)"
                ax.scatter(x_valid, y_valid, marker=marker, color=color, s=40, zorder=3, label=label)
                ax.plot(x_valid, y_valid, color=color, linewidth=1, alpha=0.4)
            ax.legend(fontsize=6, loc="best", framealpha=0.8)
        else:
            ax.text(0.5, 0.5, "이상 없음", transform=ax.transAxes,
                    ha="center", va="center", fontsize=12, color="#94a3b8")
        ax.set_title("열화 / 개선", fontsize=10, fontweight="bold", color=cs.primary)

    elif ref == "scatter_gms":
        gms_params = [(c.upper(), f"gms_{c}") for c in GMS_COLUMNS]
        _plot_multi_trend(ax, gms_params, weeks_data, periods, mpl_colors, "GMS")

    # 공통 스타일
    ax.set_xticks(range(len(periods)))
    ax.set_xticklabels(periods, fontsize=6, rotation=30, ha="right")
    ax.tick_params(axis="y", labelsize=7)
    ax.grid(axis="y", alpha=0.3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    fig.tight_layout()
    img_bytes = _fig_to_png(fig, dpi=180)
    plt.close(fig)

    _place_image(slide, img_bytes,
                 Inches(elem.left), Inches(elem.top),
                 Inches(elem.width), Inches(elem.height))


def _plot_multi_trend(ax, params: list[tuple[str, str]], weeks_data, periods,
                      colors: list[str], title: str):
    """멀티 레전드 scatter+line 트렌드."""
    if params:
        for idx, (display, data_key) in enumerate(params):
            values = _extract_values(weeks_data, data_key)
            x_valid = [i for i, v in enumerate(values) if v is not None]
            y_valid = [values[i] for i in x_valid]
            color = colors[idx % len(colors)]
            ax.scatter(x_valid, y_valid, marker="o", color=color, s=40, zorder=3, label=display)
            ax.plot(x_valid, y_valid, color=color, linewidth=1, alpha=0.5)
        ax.legend(fontsize=7, loc="best", framealpha=0.8)
    else:
        ax.text(0.5, 0.5, "이상 없음", transform=ax.transAxes,
                ha="center", va="center", fontsize=12, color="#94a3b8")
    ax.set_title(title, fontsize=10, fontweight="bold")


def _extract_values(weeks_data: list[dict], data_key: str) -> list[float | None]:
    values = []
    for wd in weeks_data:
        v = wd.get(data_key)
        try:
            values.append(float(v))
        except (TypeError, ValueError):
            values.append(None)
    return values


# ================================================================
# 유틸리티
# ================================================================
def _add_textbox(slide, text: str, left, top, width, height,
                 font_size=Pt(12), font_color=RGBColor(0x1E, 0x29, 0x3B),
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _style_cell(cell, font_size, bold, font_color, fill_color, font_name="맑은 고딕"):
    for p in cell.text_frame.paragraphs:
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def _place_image(slide, img_bytes: bytes, left, top, max_width, max_height):
    try:
        from PIL import Image as PILImage
        img = PILImage.open(io.BytesIO(img_bytes))
        img_w, img_h = img.size
        ratio = min(max_width / img_w, max_height / img_h)
        final_w = int(img_w * ratio)
        final_h = int(img_h * ratio)
        offset_x = (max_width - final_w) // 2
        offset_y = (max_height - final_h) // 2
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            left + offset_x, top + offset_y, final_w, final_h,
        )
    except Exception as e:
        logger.warning("[PPT Renderer] 이미지 삽입 실패: %s", e)


def _fig_to_png(fig, dpi: int = 150) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    buf.seek(0)
    return buf.read()


def _render_multiline_text(slide, text: str, elem: ElementDesign,
                           cs, font: str):
    """여러 줄 텍스트를 textbox에 렌더링 (줄바꿈 보존)."""
    txBox = slide.shapes.add_textbox(
        Inches(elem.left), Inches(elem.top),
        Inches(elem.width), Inches(elem.height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(elem.resolved_style.font_size or 12)
        p.font.color.rgb = _hex_to_rgb(elem.resolved_style.font_color or cs.text)
        p.font.name = font
        # 빈 줄은 간격 조절
        if not line.strip():
            p.space_after = Pt(4)


def _fmt_val(val) -> str:
    if val in (None, "-", ""):
        return "-"
    try:
        return f"{float(val):.2f}"
    except (TypeError, ValueError):
        return str(val)
