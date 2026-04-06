"""
Yield Report PPT Builder
=========================
python-pptx + matplotlib 기반 수율 분석 리포트 PPT 생성기.
HTML 의존 없이 state의 raw 데이터(weeks_data, anomaly_params 등)로
네이티브 테이블/차트를 직접 생성합니다.

슬라이드 구성:
  1. 타이틀
  2. 수율 요약 (핵심 지표 카드 + 이상 파라미터 테이블)
  3. 수율 데이터 테이블 (pt1h → weeks_data에서 직접 생성)
  4. 파라미터 트렌드 차트 (matplotlib scatter/line)
  5. GMS 수율 추이 차트 (matplotlib bar)
  6. Wafer Map (base64 PNG 이미지 — cummap/binmap)
  7. WADS 리포트 (텍스트 요약)
  8. 불량 이력 (텍스트 요약)
  9. AI 분석 결과
 10. 엔딩

사용법:
    builder = YieldReportPPTBuilder()
    pptx_bytes, fpath = builder.build(state)
"""
from __future__ import annotations

import base64
import io
import logging
import os
import re
import uuid
from datetime import date, datetime
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

# matplotlib 한글 폰트 설정 (서버 환경에 맞게 조정)
_KR_FONTS = ["Malgun Gothic", "NanumGothic", "AppleGothic", "DejaVu Sans"]
for _f in _KR_FONTS:
    try:
        matplotlib.font_manager.findfont(_f, fallback_to_default=False)
        plt.rcParams["font.family"] = _f
        break
    except Exception:
        continue
plt.rcParams["axes.unicode_minus"] = False

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from common import (
    PARA_COLUMNS, PT1C_COLUMNS, GMS_COLUMNS,
    HIGHER_IS_BETTER, GMS_HIGHER_IS_BETTER,
)

logger = logging.getLogger("yield_agent.ppt_builder")

# ── 디자인 상수 ──────────────────────────────────────────
SLIDE_WIDTH = Emu(12192000)   # 16:9 (13.33 in)
SLIDE_HEIGHT = Emu(6858000)   # 7.5 in

# 색상 팔레트
COLOR_PRIMARY    = RGBColor(0x1E, 0x29, 0x3B)
COLOR_ACCENT     = RGBColor(0x38, 0x82, 0xF6)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
COLOR_DARK_GRAY  = RGBColor(0x47, 0x55, 0x69)
COLOR_RED        = RGBColor(0xEF, 0x44, 0x44)
COLOR_GREEN      = RGBColor(0x22, 0xC5, 0x5E)
COLOR_YELLOW     = RGBColor(0xF5, 0x9E, 0x0B)

# matplotlib 색상 (hex)
MPL_BLUE    = "#3b82f6"
MPL_RED     = "#ef4444"
MPL_GREEN   = "#22c55e"
MPL_AMBER   = "#f59e0b"
MPL_PURPLE  = "#8b5cf6"
MPL_CYAN    = "#06b6d4"
MPL_COLORS  = [MPL_BLUE, MPL_RED, MPL_GREEN, MPL_AMBER, MPL_PURPLE, MPL_CYAN]

FONT_BODY = "맑은 고딕"

# 수율 테이블에 표시할 핵심 파라미터 (PPT 슬라이드 폭 제한)
KEY_PT1H_PARAMS = ["VTH", "IDSAT", "IDLIN", "IOFF", "ION", "IGATE", "IDDQ",
                   "VMIN", "FMAX", "GM_MAX", "SS", "DIBL"]

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "generated")
os.makedirs(OUTPUT_DIR, exist_ok=True)


class YieldReportPPTBuilder:
    """Yield Agent state로부터 PPTX 파일을 생성하는 빌더."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    # ── Public API ───────────────────────────────────────
    def build_compact(self, state: dict[str, Any]) -> tuple[bytes, str]:
        """compact 베이스 + 동적 추가 슬라이드 PPT 생성.

        파이프라인:
        1. state → GLM-4.7 (디자인 JSON) → ppt_renderer (베이스 3슬라이드)
        2. state에 WADS/불량이력/AI분석 데이터가 있으면 각각 1장씩 추가
        LLM 호출 실패 시 기본 디자인으로 fallback합니다.
        """
        lotcd = state.get("lotcd", "Unknown")
        ref_date = state.get("ref_date", date.today().strftime("%Y%m%d"))
        weeks_data = state.get("weeks_data", [])

        if not weeks_data:
            raise ValueError(f"[PPT compact] weeks_data가 비어있어 리포트를 생성할 수 없습니다. (lotcd={lotcd})")

        from ppt_llm_designer import generate_slide_design, generate_extra_slide
        from ppt_renderer import render_presentation

        # 1) 베이스 디자인 (title + data_overview + ending)
        design = generate_slide_design(state)

        # 2) 추가 섹션: state에 데이터가 있는 섹션만 ending 앞에 삽입
        extra_sections = [
            ("wads", state.get("wads_artifacts", [])),
            ("fail_history", state.get("fail_history_artifacts", [])),
            ("analysis", state.get("analysis_result", "")),
        ]
        extra_slides = []
        for section_type, data in extra_sections:
            if data:  # 비어있지 않은 경우만
                slide_design = generate_extra_slide(state, section_type, None)
                extra_slides.append(slide_design)

        # ending 슬라이드 앞에 추가 슬라이드 삽입
        if extra_slides:
            ending = design.slides[-1]  # 마지막 = ending
            design.slides = design.slides[:-1] + extra_slides + [ending]

        # 3) 렌더링 — extra_content를 state에 임시 주입하여 전달
        for slide_design in design.slides:
            if slide_design.extra_content:
                state["_current_extra_content"] = slide_design.extra_content
            render_presentation(self.prs,
                              type(design)(color_scheme=design.color_scheme,
                                           font_family=design.font_family,
                                           slides=[slide_design]),
                              state)
            state.pop("_current_extra_content", None)

        # 저장
        buf = io.BytesIO()
        self.prs.save(buf)
        pptx_bytes = buf.getvalue()

        fname = f"yield_report_{lotcd}_{ref_date}_{uuid.uuid4().hex[:6]}.pptx"
        fpath = os.path.join(OUTPUT_DIR, fname)
        with open(fpath, "wb") as f:
            f.write(pptx_bytes)

        logger.info("[PPT Builder] compact+동적 생성 완료: %s (%d bytes, %d slides)",
                     fname, len(pptx_bytes), len(self.prs.slides))
        return pptx_bytes, fpath

    def build(self, state: dict[str, Any]) -> tuple[bytes, str]:
        """state 전체를 받아 PPT를 생성하고 (bytes, file_path)를 반환."""
        lotcd = state.get("lotcd", "Unknown")
        ref_date = state.get("ref_date", date.today().strftime("%Y%m%d"))
        unit = state.get("unit", "weekly")
        weeks_data = state.get("weeks_data", [])
        anomaly_params = state.get("anomaly_params", [])
        filter_params = state.get("filter_params") or None

        # 1) 타이틀
        self._add_title_slide(lotcd, ref_date, unit)

        # 2) 수율 요약
        if anomaly_params or weeks_data:
            self._add_yield_summary_slide(lotcd, weeks_data, anomaly_params, unit)

        # 3) 수율 데이터 테이블 (네이티브)
        if weeks_data:
            self._add_yield_table_slides(weeks_data, lotcd, anomaly_params,
                                         filter_params, unit)

        # 4) 파라미터 트렌드 차트 (matplotlib)
        if weeks_data and len(weeks_data) >= 2:
            self._add_trend_chart_slide(weeks_data, anomaly_params, unit)

        # 5) GMS 수율 추이 차트 (matplotlib)
        if weeks_data and len(weeks_data) >= 2:
            self._add_gms_chart_slide(weeks_data, unit)

        # 6) Wafer Map — base64 PNG가 포함된 artifact에서 이미지 추출
        map_artifacts = state.get("map_artifacts", [])
        yield_artifacts = state.get("yield_artifacts", [])
        # yield_cummap도 이미지 기반이므로 포함
        for art in yield_artifacts:
            if art.get("title") == "yield_cummap":
                images = self._extract_base64_images(art.get("data", ""))
                if images:
                    self._add_image_slide("누적 패스레이트 맵", images)
        for art in map_artifacts:
            images = self._extract_base64_images(art.get("data", ""))
            if images:
                title = art.get("title", "Wafer Map")
                self._add_image_slide(f"Wafer Map: {title}", images)

        # 7) WADS — 구조화된 슬라이드
        wads_artifacts = state.get("wads_artifacts", [])
        if wads_artifacts:
            self._add_section_divider("WADS 이상 탐지", "열화 파라미터 탐지 및 WADS Score 분석")
            self._add_wads_slides(wads_artifacts, anomaly_params)

        # 8) 불량 이력 — 구조화된 슬라이드
        fh_artifacts = state.get("fail_history_artifacts", [])
        if fh_artifacts:
            self._add_section_divider("불량 이력", "Fail History 및 불량 유형별 통계")
            self._add_fail_history_slides(fh_artifacts)

        # 9) AI 분석 — 개선된 레이아웃
        analysis = state.get("analysis_result", "")
        if analysis:
            self._add_section_divider("AI 분석", "Yield Agent 종합 분석 결과")
            self._add_analysis_slide_v2(analysis, anomaly_params)

        # 10) 엔딩
        self._add_ending_slide()

        # 저장
        buf = io.BytesIO()
        self.prs.save(buf)
        pptx_bytes = buf.getvalue()

        fname = f"yield_report_{lotcd}_{ref_date}_{uuid.uuid4().hex[:6]}.pptx"
        fpath = os.path.join(OUTPUT_DIR, fname)
        with open(fpath, "wb") as f:
            f.write(pptx_bytes)

        logger.info("[PPT Builder] 생성 완료: %s (%d bytes, %d slides)",
                     fname, len(pptx_bytes), len(self.prs.slides))
        return pptx_bytes, fpath

    # ================================================================
    # Compact 슬라이드 1 — 테이블 3개 + scatter plot 3개
    # ================================================================
    def _add_compact_slide1(
        self, lotcd: str, weeks_data: list[dict],
        anomaly_params: list[dict], filter_params: list | None, unit: str,
    ):
        """상단: 통합 멀티헤더 테이블, 하단: scatter plot 3개."""
        slide = self._add_blank_slide()
        self._add_slide_header(slide, f"{lotcd} 수율 요약")

        period_label = {"weekly": "WEEK", "monthly": "MONTH", "daily": "DATE"}.get(unit, "WEEK")
        anomaly_map = {a["param"]: a for a in anomaly_params} if anomaly_params else {}

        # ── 위쪽: 통합 멀티헤더 테이블 ──
        pt1h_cols = KEY_PT1H_PARAMS[:6]
        pt1c_cols = PT1C_COLUMNS[:]
        gms_cols = GMS_COLUMNS[:]

        # 그룹 정의: (그룹명, [(display_name, data_key), ...])
        groups = [
            ("PT1H", [(c, c) for c in pt1h_cols]),
            ("PT1C", [(c, f"pt1c_{c}") for c in pt1c_cols]),
            ("GMS",  [(c.upper(), f"gms_{c}") for c in gms_cols]),
        ]

        self._add_compact_multi_header_table(
            slide, weeks_data, groups, period_label, anomaly_map,
            left=Inches(0.2), top=Inches(1.15), width=Inches(12.9),
        )

        # ── 아래쪽: scatter plot 3개 ──
        self._add_compact_scatter(slide, weeks_data, anomaly_params, unit)

    def _add_compact_multi_header_table(
        self, slide, weeks_data: list[dict],
        groups: list[tuple[str, list[tuple[str, str]]]],
        period_label: str, anomaly_map: dict,
        left, top, width,
    ):
        """멀티헤더 통합 테이블: row0=그룹헤더(merge), row1=파라미터명, row2+=데이터."""
        # 전체 파라미터 수 계산
        all_params = []  # (display_name, data_key)
        group_spans = []  # (group_name, start_col, n_cols)
        col_offset = 1  # 0번 컬럼은 period
        for gname, params in groups:
            group_spans.append((gname, col_offset, len(params)))
            all_params.extend(params)
            col_offset += len(params)

        n_cols = 1 + len(all_params)
        n_rows = 2 + len(weeks_data)  # 그룹헤더 + 파라미터헤더 + 데이터
        row_h = min(Inches(0.26), Inches(2.5) / max(n_rows, 1))

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, left, top, width, row_h * n_rows,
        )
        table = table_shape.table

        # ── row 0: 그룹 헤더 (merge) ──
        # period 셀 (row0 + row1 세로 merge)
        table.cell(0, 0).merge(table.cell(1, 0))
        table.cell(0, 0).text = period_label
        self._style_cell(table.cell(0, 0), font_size=Pt(7), bold=True,
                         font_color=COLOR_WHITE, fill_color=COLOR_PRIMARY)

        for gname, start, span in group_spans:
            if span > 1:
                table.cell(0, start).merge(table.cell(0, start + span - 1))
            table.cell(0, start).text = gname
            # 그룹별 색상
            if gname == "PT1H":
                bg = COLOR_PRIMARY
            elif gname == "PT1C":
                bg = RGBColor(0x6D, 0x28, 0xD9)  # purple
            else:
                bg = RGBColor(0x05, 0x96, 0x69)  # teal
            self._style_cell(table.cell(0, start), font_size=Pt(7), bold=True,
                             font_color=COLOR_WHITE, fill_color=bg)

        # ── row 1: 파라미터 이름 ──
        for j, (display, data_key) in enumerate(all_params):
            cell = table.cell(1, 1 + j)
            cell.text = display
            base_param = data_key.replace("pt1c_", "").replace("gms_", "")
            if base_param in anomaly_map:
                bg = COLOR_RED if anomaly_map[base_param]["direction"] == "열화" else COLOR_GREEN
            else:
                bg = COLOR_DARK_GRAY
            self._style_cell(cell, font_size=Pt(6), bold=True,
                             font_color=COLOR_WHITE, fill_color=bg)

        # ── row 2+: 데이터 ──
        last_idx = len(weeks_data) - 1
        for i, wd in enumerate(weeks_data):
            row_idx = i + 2
            bg = COLOR_WHITE if i % 2 == 0 else COLOR_LIGHT_GRAY

            table.cell(row_idx, 0).text = str(wd.get("week", "?"))
            self._style_cell(table.cell(row_idx, 0), font_size=Pt(6),
                             font_color=COLOR_PRIMARY, fill_color=bg)

            for j, (display, data_key) in enumerate(all_params):
                cell = table.cell(row_idx, 1 + j)
                cell.text = self._fmt_val(wd.get(data_key, "-"))
                base_param = data_key.replace("pt1c_", "").replace("gms_", "")
                if i == last_idx and base_param in anomaly_map:
                    a = anomaly_map[base_param]
                    fc = COLOR_RED if a["direction"] == "열화" else COLOR_GREEN
                    self._style_cell(cell, font_size=Pt(6), bold=True,
                                     font_color=fc, fill_color=bg)
                else:
                    self._style_cell(cell, font_size=Pt(6),
                                     font_color=COLOR_DARK_GRAY, fill_color=bg)

    def _add_compact_scatter(self, slide, weeks_data: list[dict],
                              anomaly_params: list[dict], unit: str):
        """Compact 슬라이드 하단: PT1H+PT1C / 열화개선 / GMS scatter plot 1×3."""
        periods = [wd.get("week", f"P{i}") for i, wd in enumerate(weeks_data)]

        fig, axes = plt.subplots(1, 3, figsize=(13, 2.8))
        fig.patch.set_facecolor("white")

        def _plot_params(ax, param_list, title):
            """파라미터 리스트를 멀티 레전드 scatter로 그리는 공통 함수."""
            if param_list:
                for idx, (param, data_key) in enumerate(param_list):
                    values = []
                    for wd in weeks_data:
                        v = wd.get(data_key)
                        try:
                            values.append(float(v))
                        except (TypeError, ValueError):
                            values.append(None)
                    x_valid = [i for i, v in enumerate(values) if v is not None]
                    y_valid = [values[i] for i in x_valid]
                    color = MPL_COLORS[idx % len(MPL_COLORS)]
                    ax.scatter(x_valid, y_valid, marker="o", color=color,
                               s=40, zorder=3, label=param)
                    ax.plot(x_valid, y_valid, color=color, linewidth=1, alpha=0.5)
                ax.legend(fontsize=7, loc="best", framealpha=0.8)
            else:
                ax.text(0.5, 0.5, "이상 없음", transform=ax.transAxes,
                        ha="center", va="center", fontsize=12, color="#94a3b8")
            ax.set_title(title, fontsize=10, fontweight="bold", color="#1e293b")
            ax.set_xticks(range(len(periods)))
            ax.set_xticklabels(periods, fontsize=6, rotation=30, ha="right")
            ax.tick_params(axis="y", labelsize=7)
            ax.grid(axis="y", alpha=0.3)
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)

        # ── subplot 1: PT1H + PT1C 이상 파라미터 트렌드 ──
        pt1h_pt1c = []
        for a in anomaly_params:
            p = a["param"]
            if p in PARA_COLUMNS:
                pt1h_pt1c.append((p, p))
            elif p in PT1C_COLUMNS:
                pt1h_pt1c.append((p, f"pt1c_{p}"))
        _plot_params(axes[0], pt1h_pt1c[:6], "PT1H / PT1C")

        # ── subplot 2: 열화/개선 파라미터 트렌드 (방향별 마커) ──
        ax = axes[1]
        if anomaly_params:
            for idx, a in enumerate(anomaly_params[:6]):
                param = a["param"]
                data_key = f"pt1c_{param}" if param in PT1C_COLUMNS else param
                values = []
                for wd in weeks_data:
                    v = wd.get(data_key)
                    try:
                        values.append(float(v))
                    except (TypeError, ValueError):
                        values.append(None)
                x_valid = [i for i, v in enumerate(values) if v is not None]
                y_valid = [values[i] for i in x_valid]
                color = MPL_RED if a.get("direction") == "열화" else MPL_GREEN
                marker = "v" if a.get("direction") == "열화" else "^"
                direction = a.get('direction', '?')
                pct = a.get('change_pct', 0.0)
                label = f"{param} ({direction} {pct:+.1f}%)"
                ax.scatter(x_valid, y_valid, marker=marker, color=color,
                           s=40, zorder=3, label=label)
                ax.plot(x_valid, y_valid, color=color, linewidth=1, alpha=0.4)
            ax.legend(fontsize=6, loc="best", framealpha=0.8)
        else:
            ax.text(0.5, 0.5, "이상 없음", transform=ax.transAxes,
                    ha="center", va="center", fontsize=12, color="#94a3b8")
        ax.set_title("열화 / 개선", fontsize=10, fontweight="bold", color="#1e293b")
        ax.set_xticks(range(len(periods)))
        ax.set_xticklabels(periods, fontsize=6, rotation=30, ha="right")
        ax.tick_params(axis="y", labelsize=7)
        ax.grid(axis="y", alpha=0.3)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

        # ── subplot 3: GMS 수율 트렌드 ──
        gms_params = [(c.upper(), f"gms_{c}") for c in GMS_COLUMNS]
        _plot_params(axes[2], gms_params, "GMS")

        fig.tight_layout()
        img_bytes = self._fig_to_png(fig, dpi=180)
        plt.close(fig)

        self._place_image(slide, img_bytes,
                          Inches(0.2), Inches(3.9), Inches(12.9), Inches(3.4))

    # ================================================================
    # 슬라이드 빌더 — 수율 데이터 테이블 (네이티브)
    # ================================================================
    def _add_yield_table_slides(
        self, weeks_data: list[dict], lotcd: str,
        anomaly_params: list[dict], filter_params: list | None, unit: str,
    ):
        """weeks_data를 python-pptx 네이티브 테이블로 변환.

        컬럼 수가 많으므로 pt1h / pt1c+gms 로 2장 분할합니다.
        filter_params가 있으면 해당 파라미터만 표시합니다.
        """
        period_label = {"weekly": "WEEK", "monthly": "MONTH", "daily": "DATE"}.get(unit, "WEEK")
        anomaly_map = {a["param"]: a for a in anomaly_params} if anomaly_params else {}

        # pt1h 파라미터 선택
        if filter_params:
            pt1h_cols = [c for c in PARA_COLUMNS if c in filter_params]
        else:
            pt1h_cols = KEY_PT1H_PARAMS  # PPT용 핵심 파라미터만

        # ── Slide 1: PT1H 테이블 ──
        self._build_data_table_slide(
            title=f"{lotcd} PT1H 수율 데이터",
            weeks_data=weeks_data,
            param_cols=pt1h_cols,
            col_prefix="",
            period_label=period_label,
            anomaly_map=anomaly_map,
            show_meta=True,
        )

        # ── Slide 2: PT1C + GMS 테이블 ──
        pt1c_gms_cols = [f"pt1c_{c}" for c in PT1C_COLUMNS] + [f"gms_{c}" for c in GMS_COLUMNS]
        display_names = {f"pt1c_{c}": f"PT1C_{c}" for c in PT1C_COLUMNS}
        display_names.update({f"gms_{c}": f"GMS_{c}" for c in GMS_COLUMNS})

        self._build_data_table_slide(
            title=f"{lotcd} PT1C / GMS 수율 데이터",
            weeks_data=weeks_data,
            param_cols=pt1c_gms_cols,
            col_prefix="",
            period_label=period_label,
            anomaly_map=anomaly_map,
            show_meta=True,
            display_names=display_names,
        )

    def _build_data_table_slide(
        self, title: str, weeks_data: list[dict], param_cols: list[str],
        col_prefix: str, period_label: str, anomaly_map: dict,
        show_meta: bool = True, display_names: dict | None = None,
    ):
        """범용 데이터 테이블 슬라이드 생성."""
        if not param_cols:
            return

        slide = self._add_blank_slide()
        self._add_slide_header(slide, title)

        meta_cols = [period_label, "LOT", "WF"] if show_meta else [period_label]
        n_meta = len(meta_cols)
        n_params = len(param_cols)
        n_cols = n_meta + n_params
        n_rows = len(weeks_data) + 1  # +1 for header

        # 테이블 사이즈 계산
        table_left = Inches(0.3)
        table_top = Inches(1.3)
        table_width = Inches(12.7)
        row_height = min(Inches(0.35), Inches(5.5) / max(n_rows, 1))
        table_height = row_height * n_rows

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, table_left, table_top, table_width, table_height,
        )
        table = table_shape.table

        # 헤더
        for j, h in enumerate(meta_cols):
            cell = table.cell(0, j)
            cell.text = h
            self._style_cell(cell, font_size=Pt(8), bold=True,
                             font_color=COLOR_WHITE, fill_color=COLOR_PRIMARY)

        for j, col in enumerate(param_cols):
            cell = table.cell(0, n_meta + j)
            display = (display_names or {}).get(col, col)
            cell.text = display
            # 이상 파라미터면 강조 색상
            base_param = col.replace("pt1c_", "").replace("gms_", "")
            if base_param in anomaly_map:
                bg = COLOR_RED if anomaly_map[base_param]["direction"] == "열화" else COLOR_GREEN
            else:
                bg = COLOR_PRIMARY
            self._style_cell(cell, font_size=Pt(7), bold=True,
                             font_color=COLOR_WHITE, fill_color=bg)

        # 데이터 행
        last_idx = len(weeks_data) - 1
        for i, wd in enumerate(weeks_data):
            row_idx = i + 1
            bg = COLOR_WHITE if i % 2 == 0 else COLOR_LIGHT_GRAY

            # 메타 컬럼
            if show_meta:
                meta_vals = [
                    str(wd.get("week", "?")),
                    str(wd.get("lotcount", "-")),
                    str(wd.get("wfCount", "-")),
                ]
            else:
                meta_vals = [str(wd.get("week", "?"))]

            for j, v in enumerate(meta_vals):
                cell = table.cell(row_idx, j)
                cell.text = v
                self._style_cell(cell, font_size=Pt(8), font_color=COLOR_PRIMARY,
                                 fill_color=bg)

            # 파라미터 값
            for j, col in enumerate(param_cols):
                cell = table.cell(row_idx, n_meta + j)
                raw_val = wd.get(col, "-")
                display = self._fmt_val(raw_val)
                cell.text = display

                # 마지막 기간 + 이상 파라미터면 색상 강조
                base_param = col.replace("pt1c_", "").replace("gms_", "")
                if i == last_idx and base_param in anomaly_map:
                    a = anomaly_map[base_param]
                    fc = COLOR_RED if a["direction"] == "열화" else COLOR_GREEN
                    self._style_cell(cell, font_size=Pt(8), bold=True,
                                     font_color=fc, fill_color=bg)
                else:
                    self._style_cell(cell, font_size=Pt(8), font_color=COLOR_DARK_GRAY,
                                     fill_color=bg)

    # ================================================================
    # 슬라이드 빌더 — 파라미터 트렌드 차트 (matplotlib)
    # ================================================================
    def _add_trend_chart_slide(self, weeks_data: list[dict],
                                anomaly_params: list[dict], unit: str):
        """핵심 파라미터 + 이상 파라미터의 기간별 트렌드를 matplotlib 차트로 생성."""
        # 표시할 파라미터 선정: 고정 2개 + 이상 top 4
        chart_params = ["VTH", "IDSAT"]
        anomaly_param_names = [a["param"] for a in anomaly_params
                               if a["param"] not in chart_params][:4]
        chart_params.extend(anomaly_param_names)

        periods = [wd.get("week", f"P{i}") for i, wd in enumerate(weeks_data)]
        n_params = len(chart_params)

        # 2행 × 3열 그리드 (최대 6개 파라미터)
        n_cols_grid = min(3, n_params)
        n_rows_grid = (n_params + n_cols_grid - 1) // n_cols_grid

        fig, axes = plt.subplots(n_rows_grid, n_cols_grid,
                                  figsize=(13, 2.5 * n_rows_grid),
                                  squeeze=False)
        fig.patch.set_facecolor("white")

        for idx, param in enumerate(chart_params):
            r, c = divmod(idx, n_cols_grid)
            ax = axes[r][c]

            # 데이터 추출 (pt1c_ prefix 처리)
            data_key = param
            if param in PT1C_COLUMNS:
                data_key = f"pt1c_{param}"

            values = []
            for wd in weeks_data:
                v = wd.get(data_key)
                try:
                    values.append(float(v))
                except (TypeError, ValueError):
                    values.append(None)

            # None 처리 — 유효한 값만 플롯
            x_valid = [i for i, v in enumerate(values) if v is not None]
            y_valid = [values[i] for i in x_valid]
            labels_valid = [periods[i] for i in x_valid]

            color = MPL_COLORS[idx % len(MPL_COLORS)]

            # 이상 파라미터면 방향 표시
            anomaly_info = next((a for a in anomaly_params if a["param"] == param), None)
            marker = "v" if (anomaly_info and anomaly_info["direction"] == "열화") else "o"

            ax.plot(x_valid, y_valid, marker=marker, color=color, linewidth=2,
                    markersize=6, markerfacecolor="white", markeredgewidth=2)
            ax.fill_between(x_valid, y_valid, alpha=0.1, color=color)

            # 스타일
            ax.set_title(param, fontsize=12, fontweight="bold", color="#1e293b")
            ax.set_xticks(range(len(periods)))
            ax.set_xticklabels(periods, fontsize=8, rotation=30, ha="right")
            ax.tick_params(axis="y", labelsize=8)
            ax.grid(axis="y", alpha=0.3)
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)

            # 이상 정보 라벨
            if anomaly_info:
                direction = anomaly_info["direction"]
                pct = anomaly_info["change_pct"]
                badge_color = "#ef4444" if direction == "열화" else "#22c55e"
                dir_en = "Degraded" if direction == "열화" else "Improved"
                ax.annotate(
                    f"{dir_en} {pct:+.1f}%",
                    xy=(x_valid[-1], y_valid[-1]),
                    fontsize=9, fontweight="bold", color=badge_color,
                    xytext=(5, 10), textcoords="offset points",
                )

        # 빈 서브플롯 숨기기
        for idx in range(n_params, n_rows_grid * n_cols_grid):
            r, c = divmod(idx, n_cols_grid)
            axes[r][c].set_visible(False)

        fig.suptitle("Parameter Trend", fontsize=16, fontweight="bold",
                      color="#1e293b", y=1.02)
        fig.tight_layout()

        img_bytes = self._fig_to_png(fig)
        plt.close(fig)

        slide = self._add_blank_slide()
        self._add_slide_header(slide, "파라미터 트렌드 차트")
        self._place_image(slide, img_bytes,
                          Inches(0.3), Inches(1.2), Inches(12.7), Inches(6))

    # ================================================================
    # 슬라이드 빌더 — GMS 수율 차트
    # ================================================================
    def _add_gms_chart_slide(self, weeks_data: list[dict], unit: str):
        """GMS 수율(cum0, cum2, fab, prb, pnt) 추이를 막대+선 차트로 생성."""
        periods = [wd.get("week", f"P{i}") for i, wd in enumerate(weeks_data)]
        gms_data = {}
        for col in GMS_COLUMNS:
            vals = []
            for wd in weeks_data:
                v = wd.get(f"gms_{col}")
                try:
                    vals.append(float(v))
                except (TypeError, ValueError):
                    vals.append(0)
            gms_data[col] = vals

        # 유효 데이터 확인
        if all(all(v == 0 for v in vals) for vals in gms_data.values()):
            return  # 전부 0이면 스킵

        fig, ax = plt.subplots(figsize=(13, 5))
        fig.patch.set_facecolor("white")

        x = np.arange(len(periods))
        width = 0.15
        offsets = np.arange(len(GMS_COLUMNS)) - len(GMS_COLUMNS) / 2 + 0.5

        for i, col in enumerate(GMS_COLUMNS):
            color = MPL_COLORS[i % len(MPL_COLORS)]
            bars = ax.bar(x + offsets[i] * width, gms_data[col], width,
                          label=col.upper(), color=color, alpha=0.85)

        ax.set_xlabel("Period", fontsize=10)
        ax.set_ylabel("Yield (%)", fontsize=10)
        ax.set_title("GMS Yield Trend", fontsize=14, fontweight="bold", color="#1e293b")
        ax.set_xticks(x)
        ax.set_xticklabels(periods, fontsize=9, rotation=30, ha="right")
        ax.legend(fontsize=9, loc="upper left", framealpha=0.9)
        ax.grid(axis="y", alpha=0.3)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        fig.tight_layout()

        img_bytes = self._fig_to_png(fig)
        plt.close(fig)

        slide = self._add_blank_slide()
        self._add_slide_header(slide, "GMS 수율 추이")
        self._place_image(slide, img_bytes,
                          Inches(0.3), Inches(1.2), Inches(12.7), Inches(6))

    # ================================================================
    # 타이틀 / 요약 / 분석 / 엔딩 슬라이드
    # ================================================================
    def _add_title_slide(self, lotcd: str, ref_date: str, unit: str):
        slide = self._add_blank_slide()
        self._fill_background(slide, COLOR_PRIMARY)

        self._add_textbox(
            slide, f"{lotcd} 수율 분석 리포트",
            left=Inches(1), top=Inches(2.2), width=Inches(11.33), height=Inches(1.2),
            font_size=Pt(40), font_color=COLOR_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        try:
            dt = datetime.strptime(ref_date, "%Y%m%d")
            date_str = dt.strftime("%Y년 %m월 %d일")
        except ValueError:
            date_str = ref_date
        unit_label = {"weekly": "주간", "monthly": "월간", "daily": "일별"}.get(unit, unit)

        self._add_textbox(
            slide,
            f"기준일: {date_str}  |  단위: {unit_label}  |  생성: {date.today().strftime('%Y-%m-%d')}",
            left=Inches(1), top=Inches(3.5), width=Inches(11.33), height=Inches(0.6),
            font_size=Pt(16), font_color=COLOR_LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    def _add_yield_summary_slide(
        self, lotcd: str, weeks_data: list[dict],
        anomaly_params: list[dict], unit: str,
    ):
        slide = self._add_blank_slide()
        self._add_slide_header(slide, f"{lotcd} 수율 요약")

        if weeks_data:
            curr = weeks_data[-1]
            cards = [
                ("LOT 수", str(curr.get("lotcount", "-")), COLOR_ACCENT),
                ("WF 수", str(curr.get("wfCount", "-")), COLOR_ACCENT),
                ("기간", curr.get("week", "-"), COLOR_DARK_GRAY),
            ]
            card_width = Inches(3.2)
            for i, (label, value, color) in enumerate(cards):
                left = Inches(0.8) + i * (card_width + Inches(0.3))
                self._add_card(slide, left, Inches(1.5), card_width, Inches(1.2),
                               label, value, color)

        if anomaly_params:
            self._add_textbox(
                slide, f"이상 감지 파라미터 ({len(anomaly_params)}개)",
                left=Inches(0.8), top=Inches(3.0), width=Inches(11), height=Inches(0.4),
                font_size=Pt(16), font_color=COLOR_RED, bold=True,
            )
            display_params = anomaly_params[:8]
            rows = len(display_params) + 1
            table_shape = slide.shapes.add_table(
                rows, 5,
                left=Inches(0.8), top=Inches(3.5),
                width=Inches(11), height=Inches(0.35 * rows),
            )
            table = table_shape.table
            for j, h in enumerate(["파라미터", "이전값", "현재값", "변화율", "방향"]):
                cell = table.cell(0, j)
                cell.text = h
                self._style_cell(cell, font_size=Pt(10), bold=True,
                                 font_color=COLOR_WHITE, fill_color=COLOR_PRIMARY)
            for i, a in enumerate(display_params, 1):
                vals = [
                    a["param"],
                    self._fmt_val(a.get("prev_val")),
                    self._fmt_val(a.get("curr_val")),
                    f"{a['change_pct']:+.1f}%",
                    a["direction"],
                ]
                for j, v in enumerate(vals):
                    cell = table.cell(i, j)
                    cell.text = v
                    fc = (COLOR_RED if a["direction"] == "열화" else COLOR_GREEN) if j == 4 else COLOR_PRIMARY
                    bg = COLOR_WHITE if i % 2 == 1 else COLOR_LIGHT_GRAY
                    self._style_cell(cell, font_size=Pt(9), font_color=fc, fill_color=bg)
        else:
            self._add_textbox(
                slide, "이상 파라미터 없음 (±10% 기준)",
                left=Inches(0.8), top=Inches(3.5), width=Inches(11), height=Inches(0.5),
                font_size=Pt(18), font_color=COLOR_GREEN, bold=True,
            )

    def _add_analysis_slide(self, analysis: str):
        """Legacy — 하위 호환용."""
        self._add_analysis_slide_v2(analysis, [])

    def _add_ending_slide(self):
        slide = self._add_blank_slide()
        self._fill_background(slide, COLOR_PRIMARY)
        self._add_textbox(
            slide, "Thank You",
            left=Inches(1), top=Inches(2.5), width=Inches(11.33), height=Inches(1),
            font_size=Pt(36), font_color=COLOR_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        self._add_textbox(
            slide, f"Generated by Yield Agent  |  {date.today().strftime('%Y-%m-%d')}",
            left=Inches(1), top=Inches(3.8), width=Inches(11.33), height=Inches(0.5),
            font_size=Pt(14), font_color=COLOR_LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    # ================================================================
    # 섹션 구분 슬라이드
    # ================================================================
    def _add_section_divider(self, title: str, subtitle: str = ""):
        """시각적 섹션 구분 슬라이드 (다크 배경 + 액센트 라인)."""
        slide = self._add_blank_slide()
        self._fill_background(slide, COLOR_PRIMARY)

        # 왼쪽 액센트 바
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(1.0), top=Inches(2.0), width=Pt(5), height=Inches(2.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT
        bar.line.fill.background()

        self._add_textbox(
            slide, title,
            left=Inches(1.5), top=Inches(2.2), width=Inches(10), height=Inches(1),
            font_size=Pt(36), font_color=COLOR_WHITE, bold=True,
        )
        if subtitle:
            self._add_textbox(
                slide, subtitle,
                left=Inches(1.5), top=Inches(3.4), width=Inches(10), height=Inches(0.6),
                font_size=Pt(16), font_color=COLOR_LIGHT_GRAY,
            )

    # ================================================================
    # WADS 슬라이드 — 구조화된 카드 + 테이블
    # ================================================================
    def _add_wads_slides(self, wads_artifacts: list[dict],
                         anomaly_params: list[dict]):
        """WADS artifact를 파싱하여 시각적 슬라이드 생성."""
        # HTML에서 구조 데이터 추출
        html = ""
        for art in wads_artifacts:
            html += art.get("data", "")

        wads_rows = self._parse_html_table(html)
        param_sections = self._parse_wads_params(html)

        # ── Slide 1: 파라미터별 상태 카드 ──
        slide = self._add_blank_slide()
        self._add_slide_header(slide, "WADS 파라미터 상태")

        if param_sections:
            cards_per_row = min(3, len(param_sections))
            card_w = Inches(3.6)
            card_h = Inches(2.2)
            gap = Inches(0.3)
            for i, ps in enumerate(param_sections[:6]):
                row, col = divmod(i, cards_per_row)
                left = Inches(0.5) + col * (card_w + gap)
                top = Inches(1.4) + row * (card_h + gap)
                self._add_wads_status_card(slide, left, top, card_w, card_h, ps)
        else:
            # 파싱 실패시 텍스트 폴백
            text = self._strip_html(html)[:2500]
            self._add_textbox(
                slide, text,
                left=Inches(0.8), top=Inches(1.5), width=Inches(11), height=Inches(5.5),
                font_size=Pt(10), font_color=COLOR_DARK_GRAY,
            )

        # ── Slide 2: WADS Score 테이블 ──
        if wads_rows:
            slide2 = self._add_blank_slide()
            self._add_slide_header(slide2, "WADS Score 종합")

            n_cols = len(wads_rows[0]) if wads_rows else 4
            n_rows = len(wads_rows)
            table_shape = slide2.shapes.add_table(
                n_rows, n_cols,
                left=Inches(1.5), top=Inches(1.5),
                width=Inches(10), height=Inches(0.4 * n_rows),
            )
            table = table_shape.table

            for i, row in enumerate(wads_rows):
                for j, cell_text in enumerate(row):
                    cell = table.cell(i, j)
                    cell.text = cell_text.strip()
                    if i == 0:
                        self._style_cell(cell, font_size=Pt(10), bold=True,
                                         font_color=COLOR_WHITE, fill_color=COLOR_PRIMARY)
                    else:
                        # 등급 색상 처리
                        fc = COLOR_PRIMARY
                        if "critical" in cell_text.lower():
                            fc = COLOR_RED
                        elif "warning" in cell_text.lower() or "경고" in cell_text:
                            fc = COLOR_YELLOW
                        elif "normal" in cell_text.lower() or "안정" in cell_text:
                            fc = COLOR_GREEN
                        bg = COLOR_WHITE if i % 2 == 1 else COLOR_LIGHT_GRAY
                        self._style_cell(cell, font_size=Pt(10), font_color=fc,
                                         fill_color=bg)

            # 권장 조치사항 영역
            actions = self._parse_action_items(html)
            if actions:
                actions_y = Inches(1.5) + Inches(0.4 * n_rows) + Inches(0.3)
                self._add_textbox(
                    slide2, "권장 조치사항",
                    left=Inches(1.5), top=actions_y, width=Inches(10), height=Inches(0.4),
                    font_size=Pt(14), font_color=COLOR_ACCENT, bold=True,
                )
                for k, action in enumerate(actions[:5]):
                    self._add_action_item(
                        slide2, f"{k+1}. {action}",
                        left=Inches(1.7),
                        top=actions_y + Inches(0.45) + k * Inches(0.35),
                    )

    def _add_wads_status_card(self, slide, left, top, width, height, info: dict):
        """WADS 파라미터 상태 카드."""
        status = info.get("status", "normal").lower()
        if "열화" in status or "critical" in status:
            accent = COLOR_RED
            badge_text = "CRITICAL"
        elif "경고" in status or "warning" in status:
            accent = COLOR_YELLOW
            badge_text = "WARNING"
        else:
            accent = COLOR_GREEN
            badge_text = "NORMAL"

        # 카드 배경
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_LIGHT_GRAY
        card.line.width = Pt(1)

        # 상단 색상 바
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Pt(5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # 파라미터 이름
        self._add_textbox(
            slide, info.get("param", "?"),
            left=left + Inches(0.15), top=top + Inches(0.15),
            width=width - Inches(0.3), height=Inches(0.4),
            font_size=Pt(16), font_color=COLOR_PRIMARY, bold=True,
        )

        # 배지
        badge_box = self._add_textbox(
            slide, badge_text,
            left=left + width - Inches(1.3), top=top + Inches(0.15),
            width=Inches(1.15), height=Inches(0.3),
            font_size=Pt(9), font_color=accent, bold=True,
            alignment=PP_ALIGN.RIGHT,
        )

        # 설명 텍스트
        desc = info.get("desc", "")[:150]
        self._add_textbox(
            slide, desc,
            left=left + Inches(0.15), top=top + Inches(0.6),
            width=width - Inches(0.3), height=height - Inches(0.8),
            font_size=Pt(9), font_color=COLOR_DARK_GRAY,
        )

    def _add_action_item(self, slide, text: str, left, top):
        """조치사항 항목."""
        self._add_textbox(
            slide, text,
            left=left, top=top, width=Inches(9.5), height=Inches(0.3),
            font_size=Pt(10), font_color=COLOR_DARK_GRAY,
        )

    # ================================================================
    # Fail History 슬라이드 — 구조화된 테이블 + 통계
    # ================================================================
    def _add_fail_history_slides(self, fh_artifacts: list[dict]):
        """Fail History artifact를 파싱하여 시각적 슬라이드 생성."""
        html = ""
        for art in fh_artifacts:
            html += art.get("data", "")

        tables = self._parse_all_html_tables(html)

        if not tables:
            # 테이블 파싱 실패 → 텍스트 폴백
            slide = self._add_blank_slide()
            self._add_slide_header(slide, "불량 이력 (Fail History)")
            text = self._strip_html(html)[:3000]
            self._add_textbox(
                slide, text,
                left=Inches(0.8), top=Inches(1.5), width=Inches(11), height=Inches(5.5),
                font_size=Pt(10), font_color=COLOR_DARK_GRAY,
            )
            return

        # ── Slide 1: 주요 불량 이력 테이블 ──
        if len(tables) >= 1:
            main_table = tables[0]
            slide = self._add_blank_slide()
            self._add_slide_header(slide, "주요 불량 이력")
            self._render_native_table(slide, main_table, top=Inches(1.4))

        # ── Slide 2: 불량 유형별 통계 + 특이사항 ──
        if len(tables) >= 2:
            stat_table = tables[1]
            slide2 = self._add_blank_slide()
            self._add_slide_header(slide2, "불량 유형별 통계")
            self._render_native_table(slide2, stat_table, top=Inches(1.4),
                                      left=Inches(1.0), width=Inches(11))

            # 특이사항 텍스트 추출
            remarks = self._extract_section_text(html, "특이")
            if remarks:
                stat_rows = len(stat_table)
                remarks_y = Inches(1.4) + Inches(0.38 * stat_rows) + Inches(0.4)
                # 특이사항 카드
                card = slide2.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=Inches(1.0), top=remarks_y,
                    width=Inches(11), height=Inches(1.8),
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)  # amber-100
                card.line.color.rgb = COLOR_YELLOW
                card.line.width = Pt(1)

                self._add_textbox(
                    slide2, "⚠ 특이 사항",
                    left=Inches(1.2), top=remarks_y + Inches(0.1),
                    width=Inches(10.5), height=Inches(0.3),
                    font_size=Pt(12), font_color=RGBColor(0x92, 0x40, 0x0E),
                    bold=True,
                )
                self._add_textbox(
                    slide2, remarks[:500],
                    left=Inches(1.2), top=remarks_y + Inches(0.45),
                    width=Inches(10.5), height=Inches(1.2),
                    font_size=Pt(10), font_color=COLOR_DARK_GRAY,
                )

    def _render_native_table(self, slide, rows: list[list[str]],
                              top=Inches(1.5), left=Inches(0.5),
                              width=Inches(12)):
        """파싱된 2D 리스트를 python-pptx 테이블로 렌더링."""
        if not rows:
            return
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        row_h = min(Inches(0.38), Inches(5) / max(n_rows, 1))

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, left, top, width, row_h * n_rows,
        )
        table = table_shape.table

        for i, row in enumerate(rows):
            for j in range(n_cols):
                cell = table.cell(i, j)
                cell_text = row[j].strip() if j < len(row) else ""
                cell.text = cell_text

                if i == 0:
                    self._style_cell(cell, font_size=Pt(9), bold=True,
                                     font_color=COLOR_WHITE, fill_color=COLOR_PRIMARY)
                else:
                    fc = COLOR_PRIMARY
                    # 색상 힌트가 있는 셀 강조
                    if any(kw in cell_text for kw in ["증가", "↗", "열화", "Fail"]):
                        fc = COLOR_RED
                    elif any(kw in cell_text for kw in ["감소", "↘", "개선"]):
                        fc = COLOR_GREEN
                    elif any(kw in cell_text for kw in ["경고", "진행중", "→"]):
                        fc = COLOR_YELLOW
                    bg = COLOR_WHITE if i % 2 == 1 else COLOR_LIGHT_GRAY
                    self._style_cell(cell, font_size=Pt(9), font_color=fc,
                                     fill_color=bg)

    # ================================================================
    # AI 분석 슬라이드 v2 — 구조화된 레이아웃
    # ================================================================
    def _add_analysis_slide_v2(self, analysis: str, anomaly_params: list[dict]):
        """AI 분석 결과를 섹션별로 나눠 시각적으로 표현."""
        clean = re.sub(r'\[SUGGESTION:.*?\]', '', analysis).strip()

        # 분석 내용을 섹션으로 분리
        sections = self._split_analysis_sections(clean)

        slide = self._add_blank_slide()
        self._add_slide_header(slide, "AI 종합 분석 결과")

        if len(sections) >= 3:
            # 3개 이상 섹션: 요약 / 원인 / 권장 조치 카드
            section_configs = [
                ("📊 분석 요약", COLOR_ACCENT, sections[0]),
                ("🔍 원인 추정", COLOR_YELLOW, sections[1] if len(sections) > 1 else ""),
                ("✅ 권장 조치", COLOR_GREEN, sections[2] if len(sections) > 2 else ""),
            ]
            card_h = Inches(1.6)
            for i, (title, accent, content) in enumerate(section_configs):
                if not content:
                    continue
                top = Inches(1.4) + i * (card_h + Inches(0.15))
                # 좌측 액센트 바
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=Inches(0.8), top=top, width=Pt(4), height=card_h,
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = accent
                bar.line.fill.background()

                # 섹션 제목
                self._add_textbox(
                    slide, title,
                    left=Inches(1.1), top=top + Inches(0.05),
                    width=Inches(10.5), height=Inches(0.3),
                    font_size=Pt(13), font_color=accent, bold=True,
                )
                # 섹션 내용
                self._add_textbox(
                    slide, content[:600],
                    left=Inches(1.1), top=top + Inches(0.35),
                    width=Inches(10.5), height=card_h - Inches(0.4),
                    font_size=Pt(10), font_color=COLOR_DARK_GRAY,
                )
        else:
            # 섹션 분리 불가 → 단일 텍스트 (기존 방식보다는 카드형)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Inches(0.8), top=Inches(1.4),
                width=Inches(11.5), height=Inches(5.5),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            card.line.fill.background()

            if len(clean) > 3000:
                clean = clean[:3000] + "\n\n... (이하 생략)"
            self._add_textbox(
                slide, clean,
                left=Inches(1.0), top=Inches(1.6), width=Inches(11), height=Inches(5.2),
                font_size=Pt(11), font_color=COLOR_PRIMARY,
            )

    # ================================================================
    # HTML 파싱 유틸리티 (WADS / Fail History 용)
    # ================================================================
    @staticmethod
    def _parse_html_table(html: str) -> list[list[str]]:
        """HTML에서 첫 번째 <table>을 2D 리스트로 파싱."""
        table_match = re.search(r'<table[^>]*>(.*?)</table>', html, re.DOTALL)
        if not table_match:
            return []
        table_html = table_match.group(1)
        rows = re.findall(r'<tr[^>]*>(.*?)</tr>', table_html, re.DOTALL)
        result = []
        for row in rows:
            cells = re.findall(r'<t[hd][^>]*>(.*?)</t[hd]>', row, re.DOTALL)
            cleaned = [re.sub(r'<[^>]+>', '', c).strip() for c in cells]
            if cleaned:
                result.append(cleaned)
        return result

    @staticmethod
    def _parse_all_html_tables(html: str) -> list[list[list[str]]]:
        """HTML에서 모든 <table>을 파싱하여 리스트의 리스트로 반환."""
        tables = []
        for table_match in re.finditer(r'<table[^>]*>(.*?)</table>', html, re.DOTALL):
            table_html = table_match.group(1)
            rows = re.findall(r'<tr[^>]*>(.*?)</tr>', table_html, re.DOTALL)
            parsed_rows = []
            for row in rows:
                cells = re.findall(r'<t[hd][^>]*>(.*?)</t[hd]>', row, re.DOTALL)
                cleaned = [re.sub(r'<[^>]+>', '', c).strip() for c in cells]
                if cleaned:
                    parsed_rows.append(cleaned)
            if parsed_rows:
                tables.append(parsed_rows)
        return tables

    @staticmethod
    def _parse_wads_params(html: str) -> list[dict]:
        """WADS HTML에서 파라미터별 상태/설명을 추출."""
        params = []
        # <h4>숫자. 파라미터명</h4> ... <p>상태: ...</p> ... <p>설명...</p> 패턴
        sections = re.split(r'<h4[^>]*>', html)
        for sec in sections[1:]:  # 첫 번째는 h4 이전
            # 파라미터 이름 추출
            title_match = re.match(r'[\d.]*\s*(.*?)</h4>', sec)
            if not title_match:
                continue
            param_title = re.sub(r'<[^>]+>', '', title_match.group(1)).strip()
            # "VTH (Threshold Voltage)" → param_name = "VTH"
            param_name = param_title.split("(")[0].split(" ")[0].strip()

            # 상태 추출 — "상태:" 뒤의 텍스트를 태그 무관하게 추출
            status = "normal"
            status_block = re.search(r'상태.*?</p>', sec, re.DOTALL)
            if status_block:
                status_text = re.sub(r'<[^>]+>', '', status_block.group(0)).strip()
                # "상태: 열화 감지" → "열화 감지"
                status_text = re.sub(r'^상태\s*[:：]\s*', '', status_text)
                if status_text:
                    status = status_text

            # 설명 추출 (첫 번째 <p> 중 상태가 아닌 것)
            desc_parts = re.findall(r'<p[^>]*>(.*?)</p>', sec, re.DOTALL)
            desc = ""
            for dp in desc_parts:
                clean_dp = re.sub(r'<[^>]+>', '', dp).strip()
                if clean_dp and "상태" not in clean_dp[:5]:
                    desc = clean_dp
                    break

            # <li> 내용도 추가
            li_items = re.findall(r'<li[^>]*>(.*?)</li>', sec, re.DOTALL)
            li_text = " | ".join(re.sub(r'<[^>]+>', '', li).strip()
                                  for li in li_items[:3])
            if li_text:
                desc = desc + "\n" + li_text if desc else li_text

            # 비-파라미터 섹션 (종합 점수, 권장 조치 등) 필터링
            skip_keywords = ["종합", "권장", "조치", "통계", "요약", "Score", "Action"]
            if any(kw in param_title for kw in skip_keywords):
                continue

            params.append({
                "param": param_name,
                "status": status,
                "desc": desc[:200],
            })
        return params

    @staticmethod
    def _parse_action_items(html: str) -> list[str]:
        """HTML에서 조치사항 <ol>/<li> 항목 추출."""
        # "권장 조치" 또는 "조치사항" 이후의 <li> 항목 찾기
        action_section = re.search(
            r'(?:권장|조치|Action).*?<ol[^>]*>(.*?)</ol>', html, re.DOTALL | re.IGNORECASE
        )
        if not action_section:
            return []
        items = re.findall(r'<li[^>]*>(.*?)</li>', action_section.group(1), re.DOTALL)
        return [re.sub(r'<[^>]+>', '', item).strip() for item in items]

    @staticmethod
    def _extract_section_text(html: str, keyword: str) -> str:
        """특정 키워드가 포함된 <h4> 이후의 <p> 텍스트를 추출."""
        pattern = rf'<h4[^>]*>[^<]*{keyword}[^<]*</h4>\s*<p[^>]*>(.*?)</p>'
        match = re.search(pattern, html, re.DOTALL | re.IGNORECASE)
        if match:
            return re.sub(r'<[^>]+>', '', match.group(1)).strip()
        return ""

    @staticmethod
    def _split_analysis_sections(text: str) -> list[str]:
        """AI 분석 텍스트를 주요 섹션으로 분리.

        '요약/분석 결과', '원인 추정', '권장 조치' 등의 키워드로 분리 시도.
        """
        # 번호 붙은 섹션 분리: "1. ...", "근본 원인...", "권장 조치..."
        parts = re.split(
            r'\n(?=(?:근본|원인|Root|Cause|권장|조치|Recommend|Action|개선))',
            text, flags=re.IGNORECASE,
        )
        if len(parts) >= 3:
            return [p.strip() for p in parts[:3]]

        # 이중 줄바꿈으로 분리
        parts = [p.strip() for p in text.split("\n\n") if p.strip()]
        if len(parts) >= 3:
            return parts[:3]

        return [text]

    # ================================================================
    # 이미지 / 텍스트 추출 슬라이드 (Wafer Map, WADS, Fail History)
    # ================================================================
    def _add_image_slide(self, title: str, images: list[bytes]):
        if not images:
            return
        chunks = [images[i:i+4] for i in range(0, len(images), 4)]
        for ci, chunk in enumerate(chunks):
            slide = self._add_blank_slide()
            self._add_slide_header(slide, title if ci == 0 else f"{title} (계속)")
            n = len(chunk)
            if n == 1:
                self._place_image(slide, chunk[0],
                                  Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.8))
            elif n == 2:
                for i, img in enumerate(chunk):
                    self._place_image(slide, img,
                                      Inches(0.5) + i * Inches(6.2), Inches(1.3),
                                      Inches(6), Inches(5.8))
            else:
                positions = [
                    (Inches(0.5), Inches(1.3), Inches(6), Inches(2.8)),
                    (Inches(6.7), Inches(1.3), Inches(6), Inches(2.8)),
                    (Inches(0.5), Inches(4.2), Inches(6), Inches(2.8)),
                    (Inches(6.7), Inches(4.2), Inches(6), Inches(2.8)),
                ]
                for i, img in enumerate(chunk):
                    self._place_image(slide, img, *positions[i])

    def _add_text_slide(self, title: str, artifacts: list[dict]):
        """HTML artifact에서 텍스트를 추출하여 슬라이드에 표시."""
        slide = self._add_blank_slide()
        self._add_slide_header(slide, title)

        all_text = []
        all_images = []
        for art in artifacts:
            data = art.get("data", "")
            if not data:
                continue
            images = self._extract_base64_images(data)
            all_images.extend(images)
            text = self._strip_html(data)
            if text:
                all_text.append(text[:1500])

        if all_images:
            # 이미지가 있으면 이미지 우선
            self._place_image(slide, all_images[0],
                              Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.8))
        elif all_text:
            combined = "\n\n".join(all_text)[:3000]
            self._add_textbox(
                slide, combined,
                left=Inches(0.8), top=Inches(1.5), width=Inches(11), height=Inches(5.5),
                font_size=Pt(10), font_color=COLOR_DARK_GRAY,
            )
        else:
            self._add_textbox(
                slide, f"{title} 데이터가 없습니다.",
                left=Inches(0.8), top=Inches(3), width=Inches(11), height=Inches(1),
                font_size=Pt(16), font_color=COLOR_DARK_GRAY,
                alignment=PP_ALIGN.CENTER,
            )

    # ================================================================
    # 유틸리티 메서드
    # ================================================================
    def _add_blank_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _fill_background(self, slide, color: RGBColor):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_slide_header(self, slide, text: str):
        self._add_textbox(
            slide, text,
            left=Inches(0.8), top=Inches(0.3), width=Inches(11), height=Inches(0.7),
            font_size=Pt(24), font_color=COLOR_PRIMARY, bold=True,
        )
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0.8), top=Inches(1.05), width=Inches(2), height=Pt(3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_ACCENT
        shape.line.fill.background()

    def _add_textbox(
        self, slide, text: str,
        left, top, width, height,
        font_size=Pt(12), font_color=COLOR_PRIMARY, bold=False,
        alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_card(self, slide, left, top, width, height,
                  label: str, value: str, accent_color: RGBColor):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE
        shape.shadow.inherit = False
        shape.line.color.rgb = COLOR_LIGHT_GRAY
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK_GRAY
        p.font.name = FONT_BODY

        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(24)
        p2.font.color.rgb = accent_color
        p2.font.bold = True
        p2.font.name = FONT_BODY

    def _style_cell(self, cell, font_size=Pt(10), bold=False,
                    font_color=COLOR_PRIMARY, fill_color=None):
        for p in cell.text_frame.paragraphs:
            p.font.size = font_size
            p.font.bold = bold
            p.font.color.rgb = font_color
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

    def _place_image(self, slide, img_bytes: bytes, left, top, max_width, max_height):
        try:
            from PIL import Image as PILImage
            img = PILImage.open(io.BytesIO(img_bytes))
            img_w, img_h = img.size
            ratio = min(max_width / img_w, max_height / img_h)
            final_w = int(img_w * ratio)
            final_h = int(img_h * ratio)
            offset_x = (max_width - final_w) // 2
            offset_y = (max_height - final_h) // 2
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                left + offset_x, top + offset_y, final_w, final_h,
            )
        except ImportError:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes), left, top, max_width, max_height,
            )
        except Exception as e:
            logger.warning("[PPT Builder] 이미지 삽입 실패: %s", e)

    @staticmethod
    def _fig_to_png(fig, dpi: int = 150) -> bytes:
        """matplotlib Figure → PNG bytes."""
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        buf.seek(0)
        return buf.read()

    @staticmethod
    def _fmt_val(val) -> str:
        if val in (None, "-", ""):
            return "-"
        try:
            return f"{float(val):.2f}"
        except (TypeError, ValueError):
            return str(val)

    @staticmethod
    def _extract_base64_images(html: str) -> list[bytes]:
        pattern = r'src=["\']data:image/(?:png|jpeg|jpg|gif);base64,([^"\']+)["\']'
        matches = re.findall(pattern, html)
        images = []
        for b64 in matches:
            try:
                images.append(base64.b64decode(b64))
            except Exception:
                continue
        return images

    @staticmethod
    def _strip_html(html: str) -> str:
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<br\s*/?>', '\n', text)
        text = re.sub(r'</?(p|div|tr|li|h[1-6])[^>]*>', '\n', text)
        text = re.sub(r'<[^>]+>', '', text)
        text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
        text = text.replace('&nbsp;', ' ').replace('&quot;', '"')
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' {2,}', ' ', text)
        return text.strip()
